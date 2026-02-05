#!/usr/bin/env python3
"""
make_slides.py
Génère un fichier PPTX de synthèse de la Phase 1 à partir d'un fichier BibTeX (references.bib).

Usage:
    pip install python-pptx bibtexparser
    python make_slides.py output.pptx

Le script crée une présentation avec :
- slide de titre
- slides pour les sections principales
- slide "Références" listant les entrées du .bib
"""
import sys
from pptx import Presentation
from pptx.util import Pt, Inches
import bibtexparser
from pathlib import Path

BIB_PATH = "references.bib"

def read_bib_entries(path):
    if not Path(path).exists():
        print(f"Fichier {path} introuvable.")
        return []
    with open(path, encoding='utf-8') as bibfile:
        bib_database = bibtexparser.load(bibfile)
    entries = []
    for e in bib_database.entries:
        author = e.get('author', '').replace(' and ', ', ')
        title = e.get('title', '')
        year = e.get('year', '')
        entries.append({'author': author, 'title': title, 'year': year})
    return entries

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass

def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

def add_references_slide(prs, references, title="Références"):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    # Split references into readable lines, max ~6 per slide may be needed. We'll fit as many as possible.
    for i, r in enumerate(references):
        line = f"{r.get('author','')} — {r.get('title','')} ({r.get('year','')})"
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

def build_presentation(out_path):
    prs = Presentation()
    add_title_slide(prs, "EcoDriveAI — Synthèse Phase 1", "Analyse du besoin & étude bibliographique")
    outline = [
        ("Contexte & objectifs", [
            "Comprendre leviers de consommation des VE",
            "Définir cas d'usage : prédiction, recommandations, eco-routing"
        ]),
        ("Facteurs influents (exemples)", [
            "Vitesse, accélérations, topographie, météo",
            "Etat batterie (SoC/SoH), HVAC, masse, style de conduite"
        ]),
        ("Approches & modèles", [
            "Physique / grey-box + ML (résidus)",
            "ML classiques (XGBoost), DL séquentiel (LSTM/Transformer)",
            "GNN pour information spatiale, Eco-routing multi-critère"
        ]),
        ("Données & enrichissements", [
            "CAN / traces GPS, profils altitude, météo, trafic",
            "Features: pente, moyennes, counts (stops, accel > thres.)"
        ]),
        ("Métriques & protocole", [
            "MAE / RMSE / MAPE (attention aux petites valeurs)",
            "Cross-validation leave-one-route / vehicle-out",
            "Comparaison eco-routing vs shortest-time/distance"
        ]),
        ("Prochaines étapes", [
            "Rassembler jeux de données accessibles",
            "Implémenter baseline physique + XGBoost",
            "Evaluer et itérer"
        ])
    ]
    for title, bullets in outline:
        add_bullets_slide(prs, title, bullets)

    # References
    refs = read_bib_entries(BIB_PATH)
    if refs:
        # chunk references to multiple slides if many
        chunk_size = 6
        for i in range(0, len(refs), chunk_size):
            chunk = refs[i:i+chunk_size]
            add_references_slide(prs, chunk, title=f"Références (suite)" if i>0 else "Références")
    else:
        add_bullets_slide(prs, "Références", ["Aucune entrée BibTeX trouvée dans references.bib"])

    prs.save(out_path)
    print(f"Présentation enregistrée dans : {out_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "synthese_phase1.pptx"
    build_presentation(out)